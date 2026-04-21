"""Document parsing and chunking."""

from dataclasses import dataclass
from pathlib import Path

from ylj.config import CHUNK_OVERLAP, CHUNK_SIZE

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".txt", ".md", ".markdown", ".mdx",
    ".csv", ".tsv",
}

# Extensions the wizard advertises-but-doesn't-index so users with legacy
# files see them in step-04 (greyed out, "convert to .docx to index")
# rather than silently wondering why their `.doc` corpus never appears in
# chat answers. Must stay disjoint from `SUPPORTED_EXTENSIONS` /
# `PARSERS.keys()` — `tests/test_filetypes_parity.py` enforces that.
UNSUPPORTED_EXTENSIONS = {".doc", ".rtf"}

# Directories to skip during recursive scanning
SKIP_DIRS = {
    "node_modules", ".git", ".venv", "venv", "__pycache__", ".tox", ".nox",
    ".mypy_cache", ".pytest_cache", ".ruff_cache", "dist", "build",
    ".eggs", "*.egg-info", ".venv-openwebui", "qdrant_data",
}


@dataclass
class Chunk:
    text: str
    source: str
    page: int | None = None
    # Normalised absolute file path the chunk came from. Unlike `source`,
    # this never carries XLSX's ``[Sheet]`` suffix — so incremental ingest
    # can delete every chunk of a modified/removed file with a single
    # filter on this key. Defaults to `source` for parsers that already
    # store a clean path there (pdf/docx/pptx/txt/md/csv).
    source_file: str | None = None

    def __post_init__(self) -> None:
        if self.source_file is None:
            self.source_file = self.source


def parse_pdf(path: Path) -> list[Chunk]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    chunks = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            chunks.append(Chunk(text=text.strip(), source=str(path), page=i + 1))
    return chunks


def parse_docx(path: Path) -> list[Chunk]:
    from docx import Document

    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if not text:
        return []
    return [Chunk(text=text, source=str(path))]


def parse_xlsx(path: Path) -> list[Chunk]:
    """Chunk a workbook by row-group, not whole-sheet.

    The old implementation concatenated every row of a sheet into one
    giant text blob that ``split_chunks`` then shredded into thousands
    of ~500-char fragments with arbitrary row boundaries. A data-heavy
    workbook could easily produce 30 k+ chunks from a single file, which
    swamps the embed/store loop and produces garbled retrieval hits.

    Instead, accumulate rows until the running length reaches
    ``CHUNK_SIZE`` chars and emit that as a single chunk. Result:
      * chunk count is bounded by content_chars / CHUNK_SIZE
      * row boundaries are preserved (no mid-row splits)
      * retrieved chunks read like coherent row groups
    ``split_chunks`` is still run over the output for the rare chunk
    that ends up longer than the target after the final row is added.
    """
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    chunks: list[Chunk] = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        buf: list[str] = []
        buf_len = 0
        source = f"{path} [{sheet}]"

        def _emit() -> None:
            nonlocal buf, buf_len
            if not buf:
                return
            chunks.append(Chunk(
                text="\n".join(buf),
                source=source,
                source_file=str(path),
            ))
            buf = []
            buf_len = 0

        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None).strip()
            if not row_text:
                continue
            buf.append(row_text)
            buf_len += len(row_text) + 1  # +1 for the joining newline
            if buf_len >= CHUNK_SIZE:
                _emit()
        _emit()
    return chunks


def parse_pptx(path: Path) -> list[Chunk]:
    from pptx import Presentation

    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        text = "\n".join(t for t in texts if t.strip())
        if text:
            chunks.append(Chunk(text=text, source=str(path), page=i + 1))
    return chunks


def parse_text(path: Path) -> list[Chunk]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    if not text.strip():
        return []
    return [Chunk(text=text.strip(), source=str(path))]


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
    # Markdown variants all parse as text. `.mdx` has JSX embedded but
    # the JSX is valid-enough text for retrieval; treating it as plain
    # text is better than skipping the file.
    ".txt": parse_text,
    ".md": parse_text,
    ".markdown": parse_text,
    ".mdx": parse_text,
    # Tabular: `.tsv` is just `.csv` with tabs; parse_text handles both
    # fine — we don't need column-aware parsing for retrieval purposes.
    ".csv": parse_text,
    ".tsv": parse_text,
}


def parse_document(path: Path) -> list[Chunk]:
    """Parse a document into raw chunks based on file type."""
    parser = PARSERS.get(path.suffix.lower())
    if parser is None:
        print(f"Skipping unsupported file: {path}")
        return []
    return parser(path)


def split_chunks(chunks: list[Chunk]) -> list[Chunk]:
    """Split raw chunks into smaller pieces for embedding."""
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )
    result = []
    for chunk in chunks:
        splits = splitter.split_text(chunk.text)
        for split in splits:
            result.append(Chunk(
                text=split,
                source=chunk.source,
                page=chunk.page,
                source_file=chunk.source_file,
            ))
    return result


def _should_skip(path: Path) -> bool:
    """Check if any parent directory is in the skip list."""
    return any(part in SKIP_DIRS or part.endswith(".egg-info") for part in path.parts)


def load_documents(directory: Path) -> list[Chunk]:
    """Load and chunk all supported documents from a directory."""
    all_chunks = []
    for path in sorted(directory.rglob("*")):
        if (
            path.is_file()
            and path.suffix.lower() in SUPPORTED_EXTENSIONS
            and not _should_skip(path)
        ):
            print(f"Processing: {path}")
            raw = parse_document(path)
            chunked = split_chunks(raw)
            all_chunks.extend(chunked)
            print(f"  -> {len(chunked)} chunks")
    print(f"Total chunks: {len(all_chunks)}")
    return all_chunks
